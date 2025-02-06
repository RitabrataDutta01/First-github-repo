import os
import tkinter as tk
from tkinter import filedialog as fd
from pptx import Presentation
from fpdf import FPDF

root = tk.Tk()

root.title("PPTX to PDF")

root.geometry("670x360")

fn1 = " "
fn1_label = tk.Label(root, text="")
fn1_label.grid(row=0, column=2)

def oc1():
    global fn1
    fn1 = fd.askopenfilename(title="Select PowerPoint File", filetypes=[("PowerPoint Presentation", "*.pptx")])
    fn1_label.config(text=os.path.basename(fn1))

def clean_text(text):
    text = text.replace('’', "'").replace('“', '"').replace('”', '"')
    return text

def oc(fn1):
    if not fn1:
        return

    prs = Presentation(fn1)
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    
    for i, slide in enumerate(prs.slides):
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):  
                cleaned_text = clean_text(shape.text)
                
                lines = cleaned_text.split('\n')
                
                for line in lines:
                    pdf.multi_cell(0, 10, line)
                
        pdf.ln(10)

    output_pdf = os.path.join(os.path.dirname(fn1), "converted.pdf")
    pdf.output(output_pdf)
    status_label.config(text=f"PDF saved as {output_pdf}")

lbl1 = tk.Label(root, text="Choose the PowerPoint file:")
btn1 = tk.Button(root, text="Choose", command=oc1)

lbl1.grid(row=0, column=0)
btn1.grid(row=0, column=1)

lbl = tk.Label(root, text="\n")
lbl.grid(row=1, column=0)

lbl3 = tk.Label(root, text="The converted PDF will be present in the same directory as the original PowerPoint file with the name 'converted.pdf'")
lbl3.grid(row=2, column=0)

lbl4 = tk.Label(root, text="\n")
lbl4.grid(row=3, column=0)

lbl2 = tk.Label(root, text="Click to convert to PDF")
btn2 = tk.Button(root, text="Convert", command=lambda: oc(fn1))

lbl2.grid(row=4, column=0)
btn2.grid(row=4, column=1)

status_label = tk.Label(root, text="")
status_label.grid(row=5, column=1)

root.mainloop()
