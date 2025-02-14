import os, sys, threading, webbrowser
from pptx import Presentation
from fpdf import FPDF
from flask import Flask, render_template, request, send_file, url_for, redirect
from PyQt5.QtCore import QUrl
from PyQt5.QtWidgets import QApplication, QMainWindow, QVBoxLayout, QWidget
from PyQt5.QtWebEngineWidgets import QWebEngineView, QWebEngineProfile
from werkzeug.utils import secure_filename

app = Flask(__name__)

app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['OUTPUT_FOLDER'] = 'output'
AllowedExtension = ['pptx']

os.makedirs(app.config["UPLOAD_FOLDER"], exist_ok = True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok = True)

def get_output_folder():
    if getattr(sys, 'frozen', False):  
        base_path = sys._MEIPASS 
    else:
        base_path = os.getcwd()
    output_folder = os.path.join(base_path, "output")
    os.makedirs(output_folder, exist_ok=True) 
    return output_folder

def allowedfilename(filename):
    return '.' in filename and filename.rsplit('.',1)[1].lower() in AllowedExtension

@app.route('/')
def index():
    return render_template('PPT2PDF.html')

def clean_text(text):
    text = text.replace('’', "'").replace('“', '"').replace('”', '"')
    return text

@app.route('/convert', methods = ['POST'])

def p2p():
    if 'file' not in request.files:
        return redirect(request.url)
    
    file = request.files['file']
    
    if file.filename.strip() == '':
        return redirect(request.url)
    
    if file and allowedfilename(file.filename):
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)
        
        pp = Presentation(filepath)
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
    
        for i, slide in enumerate(pp.slides):
            pdf.add_page()
            pdf.set_font("Arial", size=12)
        
            for shape in slide.shapes:
                if hasattr(shape, "text"):  
                    cleaned_text = clean_text(shape.text)
                
                    lines = cleaned_text.split('\n')
                
                    for line in lines:
                        pdf.multi_cell(0, 10, line)
                
            pdf.ln(10)
        
        output_folder = os.path.abspath(app.config['OUTPUT_FOLDER'])
        os.makedirs(output_folder, exist_ok=True)

        output_filepath = os.path.join(output_folder, filename.rsplit('.', 1)[0] + '.pdf')
        pdf.output(output_filepath)

        return send_file(output_filepath, as_attachment=True)
    return render_template('PPT2PDF.html', message='Invalid file format, please upload a Powerpoint presentation')

def run_flask_app():
    app.run(debug = True, use_reloader = False)

class BrowserWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Powerpoint to PDF Converter")
        self.setGeometry(100, 100, 800, 600)

        layout = QVBoxLayout()
        self.browser = QWebEngineView()
        layout.addWidget(self.browser)

        container = QWidget()
        container.setLayout(layout)
        self.setCentralWidget(container)

        self.browser.setUrl(QUrl("http://127.0.0.1:5000/"))

    def closeEvent(self, event):

        os._exit(0)


def start_application():

    flask_thread = threading.Thread(target=run_flask_app)
    flask_thread.start()

    app = QApplication(sys.argv)
    window = BrowserWindow()
    window.show()
    sys.exit(app.exec_())

if __name__ == '__main__':
    start_application()

